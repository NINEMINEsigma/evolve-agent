"""Document generation tools — write docx, xlsx, pptx, pdf files.

docx: uses Node.js + docx npm package (pnpm i -g docx)
xlsx: uses openpyxl (already in requirements.txt)
pptx: uses python-pptx (pip install python-pptx)
pdf: uses fpdf2 (pip install fpdf2)

Module-import-time registration via ``registry.register()``.
"""

from __future__ import annotations

import json
import logging
import subprocess
import sys
import uuid
from pathlib import Path
from typing import Any, Dict, List

from abstract.tools.registry import registry, tool_error, tool_result
from entity.puretype import ToolDangerLevel
from component.tools.filesystem import _s as _get_sandbox

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _check_node_docx() -> str | None:
    """Check if Node.js and the 'docx' package are available."""
    # Check Node
    try:
        subprocess.run(["node", "--version"], capture_output=True, text=True, timeout=10)
    except FileNotFoundError:
        return "Node.js is required to generate docx documents. Install Node.js then run:\n  pnpm i -g docx"
    # Check docx package
    try:
        subprocess.run(
            ["node", "-e", "require('docx')"],
            capture_output=True, text=True, timeout=10,
        )
    except FileNotFoundError:
        return "Node.js is required to generate docx documents. Install Node.js then run:\n  pnpm i -g docx"
    except Exception as exc:
        return f"docx npm package check failed: {exc}\n  pnpm i -g docx"
    return None


def _check_python_pptx() -> str | None:
    """Check if python-pptx is available."""
    try:
        import pptx  # noqa: F401
        return None
    except ImportError:
        return "python-pptx library is required:\n  pip install python-pptx"


def _check_fpdf2() -> str | None:
    """Check if fpdf2 is available."""
    try:
        import fpdf  # noqa: F401
        return None
    except ImportError:
        return "fpdf2 library is required:\n  pip install fpdf2"



def _make_output_path(subdir: str, ext: str) -> tuple[Path, str]:
    """Create a unique output path under ws:{subdir}/.

    Returns (real_fs_path, http_url_segment like "subdir/uuid.ext").
    Uses the Sandbox to resolve ws: paths to the agentspace directory
    so the /uploads/ route can serve the file.
    """
    sb = _get_sandbox()
    if sb is not None:
        try:
            agentspace = sb.agentspace
            out_dir = agentspace / subdir
            out_dir.mkdir(parents=True, exist_ok=True)
            name = f"{uuid.uuid4().hex[:12]}.{ext}"
            return out_dir / name, f"{subdir}/{name}"
        except Exception as exc:
            logger.warning("Failed to create output dir in agentspace/%s: %s", subdir, exc, exc_info=True)
    # Fallback: cwd
    out_dir = Path.cwd() / subdir
    out_dir.mkdir(parents=True, exist_ok=True)
    name = f"{uuid.uuid4().hex[:12]}.{ext}"
    return out_dir / name, f"{subdir}/{name}"


# ---------------------------------------------------------------------------
# write_docx — Node.js approach
# ---------------------------------------------------------------------------


_DOCX_JS_TEMPLATE = r"""const fs = require("fs");
const {{
  Document, Packer, Paragraph, TextRun, Table, TableRow, TableCell,
  Header, Footer, AlignmentType, LevelFormat,
  TableOfContents, HeadingLevel, BorderStyle, WidthType, ShadingType,
  PageNumber, PageBreak, ExternalHyperlink,
}} = require("docx");

const content = {{CONTENT_JSON}};
const outputPath = "{{OUTPUT_PATH}}";

// --- helpers ---
const border = {{ style: BorderStyle.SINGLE, size: 1, color: "CCCCCC" }};
const borders = {{ top: border, bottom: border, left: border, right: border }};

function makeParagraph(item) {{
  if (item.type === "paragraph") {{
    return new Paragraph({{
      spacing: {{ after: 120 }},
      children: [new TextRun(item.text || "")],
    }});
  }}
  if (item.type === "heading") {{
    const level = Math.min(Math.max((item.level || 1), 1), 6);
    const headingKey = "HEADING_" + level;
    return new Paragraph({{
      heading: HeadingLevel[headingKey],
      spacing: {{ before: 240, after: 120 }},
      children: [new TextRun({{ text: item.text || "", bold: true, size: 32 - level * 2 }})],
    }});
  }}
  if (item.type === "bullet_list") {{
    return (item.items || []).map(text => new Paragraph({{
      numbering: {{ reference: "bullets", level: 0 }},
      spacing: {{ after: 60 }},
      children: [new TextRun(text)],
    }}));
  }}
  if (item.type === "numbered_list") {{
    return (item.items || []).map(text => new Paragraph({{
      numbering: {{ reference: "numbers", level: 0 }},
      spacing: {{ after: 60 }},
      children: [new TextRun(text)],
    }}));
  }}
  if (item.type === "table") {{
    const headers = (item.headers || []).map(h => new TableCell({{
      borders,
      width: {{ size: 9360 / (item.headers?.length || 1), type: WidthType.DXA }},
      shading: {{ fill: "D5E8F0", type: ShadingType.CLEAR }},
      margins: {{ top: 80, bottom: 80, left: 120, right: 120 }},
      children: [new Paragraph({{ children: [new TextRun({{ text: h, bold: true }})] }})],
    }}));
    const rows = (item.rows || []).map(row => new TableRow({{
      children: row.map(cell => new TableCell({{
        borders,
        width: {{ size: 9360 / (item.headers?.length || 1), type: WidthType.DXA }},
        margins: {{ top: 60, bottom: 60, left: 120, right: 120 }},
        children: [new Paragraph({{ children: [new TextRun(String(cell))] }})],
      }})),
    }}));
    return new Table({{
      width: {{ size: 9360, type: WidthType.DXA }},
      columnWidths: new Array(item.headers?.length || 1).fill(9360 / (item.headers?.length || 1)),
      rows: [new TableRow({{ children: headers }}), ...rows],
    }});
  }}
  return [];
}}

function flatten(arr) {{
  return arr.reduce((acc, val) => acc.concat(Array.isArray(val) ? val : [val]), []);
}}

const children = flatten((content.sections || []).map(section => {{
  const items = [];
  if (section.heading) {{
    items.push(makeParagraph({{ type: "heading", text: section.heading, level: section.level || 1 }}));
  }}
  if (section.content) {{
    items.push(...flatten(section.content.map(makeParagraph)));
  }}
  return items;
}}));

const doc = new Document({{
  styles: {{
    default: {{ document: {{ run: {{ font: "Arial", size: 24 }} }} }},
    paragraphStyles: [
      {{ id: "Heading1", name: "Heading 1", basedOn: "Normal", next: "Normal", quickFormat: true,
        run: {{ size: 32, bold: true, font: "Arial" }},
        paragraph: {{ spacing: {{ before: 240, after: 240 }}, outlineLevel: 0 }} }},
      {{ id: "Heading2", name: "Heading 2", basedOn: "Normal", next: "Normal", quickFormat: true,
        run: {{ size: 28, bold: true, font: "Arial" }},
        paragraph: {{ spacing: {{ before: 180, after: 180 }}, outlineLevel: 1 }} }},
      {{ id: "Heading3", name: "Heading 3", basedOn: "Normal", next: "Normal", quickFormat: true,
        run: {{ size: 24, bold: true, font: "Arial" }},
        paragraph: {{ spacing: {{ before: 120, after: 120 }}, outlineLevel: 2 }} }},
    ],
  }},
  numbering: {{
    config: [
      {{ reference: "bullets",
        levels: [{{ level: 0, format: LevelFormat.BULLET, text: "\u2022", alignment: AlignmentType.LEFT,
          style: {{ paragraph: {{ indent: {{ left: 720, hanging: 360 }} }} }} }}] }},
      {{ reference: "numbers",
        levels: [{{ level: 0, format: LevelFormat.DECIMAL, text: "%1.", alignment: AlignmentType.LEFT,
          style: {{ paragraph: {{ indent: {{ left: 720, hanging: 360 }} }} }} }}] }},
    ],
  }},
  sections: [{{
    properties: {{
      page: {{
        size: {{ width: 12240, height: 15840 }},
        margin: {{ top: 1440, right: 1440, bottom: 1440, left: 1440 }},
      }},
    }},
    headers: {{
      default: new Header({{
        children: [new Paragraph({{ alignment: AlignmentType.RIGHT,
          children: [new TextRun({{ text: content.title || "", size: 18, color: "888888" }})] }})],
      }}),
    }},
    footers: {{
      default: new Footer({{
        children: [new Paragraph({{ alignment: AlignmentType.CENTER,
          children: [new TextRun("Page "), new TextRun({{ children: [PageNumber.CURRENT] }})] }})],
      }}),
    }},
    children: [
      ...(content.title ? [new Paragraph({{
        alignment: AlignmentType.CENTER,
        spacing: {{ after: 400 }},
        children: [new TextRun({{ text: content.title, bold: true, size: 48, font: "Arial" }})],
      }})] : []),
      ...(content.subtitle ? [new Paragraph({{
        alignment: AlignmentType.CENTER,
        spacing: {{ after: 600 }},
        children: [new TextRun({{ text: content.subtitle, size: 28, color: "555555" }})],
      }})] : []),
      ...children,
    ],
  }}],
}});

Packer.toBuffer(doc).then(buffer => {{
  fs.writeFileSync(outputPath, buffer);
  console.log("OK:" + outputPath);
}}).catch(err => {{
  console.error("ERROR:" + err.message);
  process.exit(1);
}});
"""


def _handle_write_docx(args: dict[str, Any]) -> dict:
    """Generate a .docx file using Node.js + docx npm package."""
    dep_err = _check_node_docx()
    if dep_err:
        return tool_error(dep_err)

    title: str = str(args.get("title", "")).strip()
    subtitle: str = str(args.get("subtitle", "")).strip()
    sections: list = args.get("sections", []) or []

    content = {"title": title, "subtitle": subtitle, "sections": sections}
    content_json = json.dumps(content, ensure_ascii=False)

    output_fs, output_rel = _make_output_path("documents", "docx")

    # Generate JS file
    js_content = (
        _DOCX_JS_TEMPLATE
        .replace("{{CONTENT_JSON}}", content_json)
        .replace("{{OUTPUT_PATH}}", str(output_fs))
    )
    js_path = output_fs.with_suffix(".js")
    js_path.write_text(js_content, encoding="utf-8")

    try:
        proc = subprocess.run(
            ["node", str(js_path)],
            capture_output=True, text=True, timeout=30,
        )
    except subprocess.TimeoutExpired:
        js_path.unlink(missing_ok=True)
        return tool_error("Node.js execution timed out (30s)")
    except Exception as exc:
        js_path.unlink(missing_ok=True)
        return tool_error(f"Node.js execution failed: {exc}")

    js_path.unlink(missing_ok=True)

    if proc.returncode != 0:
        err = (proc.stderr or "").strip()
        return tool_error(f"docx generation failed: {err}")

    return tool_result(
        path=f"ws:documents/{output_rel.split('/')[-1]}",
        title=title,
        sections=len(sections),
        message=f"docx document generated: {output_rel}",
    )


# ---------------------------------------------------------------------------
# write_xlsx — openpyxl
# ---------------------------------------------------------------------------


def _handle_write_xlsx(args: dict[str, Any]) -> dict:
    """Generate an .xlsx file using openpyxl."""
    try:
        import openpyxl
        from openpyxl.styles import Font
        from openpyxl.utils import get_column_letter
    except ImportError:
        return tool_error("openpyxl library is required (already in requirements.txt)")

    filename: str = str(args.get("filename", "workbook")).strip()
    sheets: list = args.get("sheets", []) or []

    wb = openpyxl.Workbook()
    # Remove default sheet
    if wb.active is not None:
        wb.remove(wb.active)

    for sheet_def in sheets:
        sheet_name: str = str(sheet_def.get("name", "Sheet")).strip()[:31]
        ws = wb.create_sheet(title=sheet_name)

        headers: list = sheet_def.get("headers", []) or []
        rows: list = sheet_def.get("rows", []) or []

        # Write headers
        if headers:
            for col_idx, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col_idx, value=str(h))
                cell.font = Font(bold=True)
            start_row = 2
        else:
            start_row = 1

        # Write data
        for row_idx, row_data in enumerate(rows, start_row):
            for col_idx, val in enumerate(row_data, 1):
                ws.cell(row=row_idx, column=col_idx, value=val)

        # Auto-adjust column width (basic)
        for col_idx in range(1, max(len(headers or []), 1) + 1):
            max_len = len(str(headers[col_idx - 1])) if headers else 10
            for row_idx in range(1, min(len(rows) + 2, 50)):
                cell_val = ws.cell(row=row_idx, column=col_idx).value
                if cell_val:
                    max_len = max(max_len, len(str(cell_val)))
            ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len + 4, 60)

    output_fs, output_rel = _make_output_path("documents", "xlsx")
    try:
        wb.save(str(output_fs))
    except Exception as exc:
        return tool_error(f"xlsx save failed: {exc}")

    return tool_result(
        path=f"ws:documents/{output_rel.split('/')[-1]}",
        sheets=len(sheets),
        message=f"xlsx workbook generated: {output_rel}",
    )


# ---------------------------------------------------------------------------
# write_pptx — python-pptx
# ---------------------------------------------------------------------------


def _handle_write_pptx(args: dict[str, Any]) -> dict:
    """Generate a .pptx file using python-pptx."""
    dep_err = _check_python_pptx()
    if dep_err:
        return tool_error(dep_err)

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return tool_error("python-pptx library is required: pip install python-pptx")

    slides_data: list = args.get("slides", []) or []

    prs = Presentation()

    for slide_def in slides_data:
        title_text: str = str(slide_def.get("title", "")).strip()
        body_items: list = slide_def.get("body", []) or []

        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        # Title
        if title_text and slide.shapes.title:
            slide.shapes.title.text = title_text

        # Body
        if body_items and slide.placeholders:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame  # type: ignore[attr-defined]
            tf.clear()
            for i, item in enumerate(body_items):
                if isinstance(item, str):
                    if i == 0:
                        p = tf.paragraphs[0]
                        p.text = item
                    else:
                        p = tf.add_paragraph()
                        p.text = item

    output_fs, output_rel = _make_output_path("documents", "pptx")
    try:
        prs.save(str(output_fs))
    except Exception as exc:
        return tool_error(f"pptx save failed: {exc}")

    return tool_result(
        path=f"ws:documents/{output_rel.split('/')[-1]}",
        slides=len(slides_data),
        message=f"pptx presentation generated: {output_rel}",
    )


# ---------------------------------------------------------------------------
# write_pdf — fpdf2
# ---------------------------------------------------------------------------


def _handle_write_pdf(args: dict[str, Any]) -> dict:
    """Generate a PDF file using fpdf2."""
    dep_err = _check_fpdf2()
    if dep_err:
        return tool_error(dep_err)

    from fpdf import FPDF

    title: str = str(args.get("title", "")).strip()
    content: str = str(args.get("content", "")).strip()
    sections: list = args.get("sections", []) or []

    pdf = FPDF()
    pdf.add_page()

    # Title
    if title:
        pdf.set_font("Helvetica", style="B", size=20)
        pdf.cell(0, 15, title, new_x="LMARGIN", new_y="NEXT", align="C")
        pdf.ln(5)

    # Body content
    if content:
        pdf.set_font("Helvetica", size=11)
        pdf.multi_cell(0, 5.5, content)
        pdf.ln(3)

    # Sections
    for section in sections:
        heading: str = str(section.get("heading", "")).strip()
        body: str = str(section.get("body", "")).strip()

        if heading:
            pdf.set_font("Helvetica", style="B", size=14)
            pdf.cell(0, 10, heading, new_x="LMARGIN", new_y="NEXT")
        if body:
            pdf.set_font("Helvetica", size=11)
            pdf.multi_cell(0, 5.5, body)
            pdf.ln(2)

    output_fs, output_rel = _make_output_path("documents", "pdf")
    try:
        pdf.output(str(output_fs))
    except Exception as exc:
        return tool_error(f"PDF save failed: {exc}")

    return tool_result(
        path=f"ws:documents/{output_rel.split('/')[-1]}",
        title=title,
        pages=pdf.pages_count,
        message=f"PDF generated: {output_rel}",
    )


# ---------------------------------------------------------------------------
# Registration
# ---------------------------------------------------------------------------


registry.register(
    name="write_docx",
    toolset="document",
    schema={
        # 生成 Word (.docx) 文件。
        #
        # ## 前置条件
        # 必须安装 Node.js 和 docx npm 包（pnpm i -g docx）。
        # sections 参数必须非空。
        #
        # ## 调用效果
        # 根据 title、subtitle 和 sections 生成 docx 文件，保存到 ws:documents/。
        # 每个 section 可包含 heading、level(1-3) 和 content（段落、列表、表格等）。
        #
        # ## 返回
        # ```json
        # {"path": "ws:documents/xxx.docx", "title": "...", "sections": 3, "message": "docx document generated: documents/xxx.docx"}
        # ```
        #
        # ## 何时使用
        # - 生成结构化 Word 报告或文档。
        # - 将表格、列表、标题组合成文档。
        #
        # ## 副作用/注意
        # - 会写入新文件到工作空间。
        # - 依赖 Node.js 执行环境。
        "description": """Generate a Word (.docx) file.

## Prerequisites
Node.js and the docx npm package must be installed (pnpm i -g docx). The sections parameter must be non-empty.

## Effect
Generates a docx file from title, subtitle, and sections, saved under ws:documents/. Each section may contain heading, level (1-3), and content blocks (paragraphs, lists, tables, etc.).

## Returns
```json
{"path": "ws:documents/xxx.docx", "title": "...", "sections": 3, "message": "docx document generated: documents/xxx.docx"}
```

## When to Use
- Generate structured Word reports or documents.
- Combine tables, lists, and headings into a document.

## Side Effects / Notes
- Writes a new file to the workspace.
- Depends on the Node.js runtime.""",
        "parameters": {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    # 文档标题（居中、大号加粗）。
                    "description": """Document title (centered, large, bold).""",
                },
                "subtitle": {
                    "type": "string",
                    # 可选的副标题。
                    "description": """Optional subtitle.""",
                },
                "sections": {
                    "type": "array",
                    # 文档章节列表。每项包含 heading（章节标题）、level（1-3）、content（内容块数组）。
                    # 内容块类型：paragraph（文本）、bullet_list（无序列表）、numbered_list（编号列表）、table（表格，含 headers 和 rows）。
                    "description": """List of document sections. Each contains heading, level (1-3), and content (array of content blocks). Content block types: paragraph, bullet_list, numbered_list, table (with headers and rows).""",
                    "items": {
                        "type": "object",
                        "properties": {
                            "heading": {"type": "string"},
                            "level": {"type": "integer"},
                            "content": {"type": "array"},
                        },
                    },
                },
            },
            "required": ["sections"],
        },
    },
    handler=_handle_write_docx,
    emoji="📝",
    danger_level=ToolDangerLevel.write,
)

registry.register(
    name="write_xlsx",
    toolset="document",
    schema={
        # 生成 Excel (.xlsx) 工作簿。
        #
        # ## 前置条件
        # 必须安装 openpyxl（已在 requirements.txt 中）。
        # sheets 参数必须非空。
        #
        # ## 调用效果
        # 根据 sheets 定义生成 xlsx 文件，保存到 ws:documents/。
        # 每个 sheet 包含 name、headers 和 rows；表头会加粗显示。
        #
        # ## 返回
        # ```json
        # {"path": "ws:documents/xxx.xlsx", "sheets": 2, "message": "xlsx workbook generated: documents/xxx.xlsx"}
        # ```
        #
        # ## 何时使用
        # - 将表格数据导出为 Excel。
        # - 生成包含多个工作表的报表。
        #
        # ## 副作用/注意
        # - 会写入新文件到工作空间，可能覆盖同名文件。
        "description": """Generate an Excel (.xlsx) workbook.

## Prerequisites
openpyxl must be installed (listed in requirements.txt). The sheets parameter must be non-empty.

## Effect
Generates an xlsx file from the sheets definition, saved under ws:documents/. Each sheet contains name, headers, and rows; headers are displayed in bold.

## Returns
```json
{"path": "ws:documents/xxx.xlsx", "sheets": 2, "message": "xlsx workbook generated: documents/xxx.xlsx"}
```

## When to Use
- Export tabular data to Excel.
- Generate reports with multiple worksheets.

## Side Effects / Notes
- Writes a new file to the workspace and may overwrite an existing file with the same name.""",
        "parameters": {
            "type": "object",
            "properties": {
                "filename": {
                    "type": "string",
                    # 工作簿名称（不含扩展名）。
                    "description": """Workbook name without extension.""",
                },
                "sheets": {
                    "type": "array",
                    # 工作表列表。每项包含 name、headers（表头）、rows（数据行）。
                    "description": """List of sheets. Each contains name, headers, and rows.""",
                    "items": {
                        "type": "object",
                        "properties": {
                            "name": {"type": "string"},
                            "headers": {"type": "array", "items": {"type": "string"}},
                            "rows": {"type": "array", "items": {"type": "array"}},
                        },
                    },
                },
            },
            "required": ["sheets"],
        },
    },
    handler=_handle_write_xlsx,
    emoji="📊",
    danger_level=ToolDangerLevel.write,
)

registry.register(
    name="write_pptx",
    toolset="document",
    schema={
        # 生成 PowerPoint (.pptx) 演示文稿。
        #
        # ## 前置条件
        # 必须安装 python-pptx（pip install python-pptx）。
        # slides 参数必须非空。
        #
        # ## 调用效果
        # 根据 slides 列表生成 pptx 文件，保存到 ws:documents/。
        # 每张幻灯片包含 title 和 body（文本行数组）。
        #
        # ## 返回
        # ```json
        # {"path": "ws:documents/xxx.pptx", "slides": 5, "message": "pptx presentation generated: documents/xxx.pptx"}
        # ```
        #
        # ## 何时使用
        # - 将大纲内容生成 PPT。
        # - 快速生成多页演示文稿。
        #
        # ## 副作用/注意
        # - 会写入新文件到工作空间，可能覆盖同名文件。
        # - 布局固定为“标题和内容”。
        "description": """Generate a PowerPoint (.pptx) presentation.

## Prerequisites
python-pptx must be installed (pip install python-pptx). The slides parameter must be non-empty.

## Effect
Generates a pptx file from the slides list, saved under ws:documents/. Each slide contains a title and a body array of text lines.

## Returns
```json
{"path": "ws:documents/xxx.pptx", "slides": 5, "message": "pptx presentation generated: documents/xxx.pptx"}
```

## When to Use
- Convert an outline into a presentation.
- Quickly generate a multi-slide deck.

## Side Effects / Notes
- Writes a new file to the workspace and may overwrite an existing file with the same name.
- Layout is fixed to "Title and Content".""",
        "parameters": {
            "type": "object",
            "properties": {
                "slides": {
                    "type": "array",
                    # 幻灯片列表。每项包含 title 和 body（文本行数组）。
                    "description": """List of slides. Each contains title and body (array of text lines).""",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "body": {"type": "array", "items": {"type": "string"}},
                        },
                    },
                },
            },
            "required": ["slides"],
        },
    },
    handler=_handle_write_pptx,
    emoji="📽️",
    danger_level=ToolDangerLevel.write,
)

registry.register(
    name="write_pdf",
    toolset="document",
    schema={
        # 生成 PDF 文件。
        #
        # ## 前置条件
        # 必须安装 fpdf2（pip install fpdf2）。
        #
        # ## 调用效果
        # 根据 title、content 和 sections 生成 PDF，保存到 ws:documents/。
        # 默认使用 Helvetica 字体。
        #
        # ## 返回
        # ```json
        # {"path": "ws:documents/xxx.pdf", "title": "...", "pages": 2, "message": "PDF generated: documents/xxx.pdf"}
        # ```
        #
        # ## 何时使用
        # - 生成简单 PDF 报告或文档。
        # - 将纯文本内容快速导出为 PDF。
        #
        # ## 副作用/注意
        # - 会写入新文件到工作空间，可能覆盖同名文件。
        # - 仅支持基础排版，复杂格式请使用 write_docx 后再转换。
        "description": """Generate a PDF file.

## Prerequisites
fpdf2 must be installed (pip install fpdf2).

## Effect
Generates a PDF from title, content, and sections, saved under ws:documents/. Uses Helvetica font by default.

## Returns
```json
{"path": "ws:documents/xxx.pdf", "title": "...", "pages": 2, "message": "PDF generated: documents/xxx.pdf"}
```

## When to Use
- Generate simple PDF reports or documents.
- Quickly export plain text content to PDF.

## Side Effects / Notes
- Writes a new file to the workspace and may overwrite an existing file with the same name.
- Supports only basic layout; for complex formatting use write_docx and convert.""",
        "parameters": {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    # PDF 标题。默认字体 Helvetica。
                    "description": """PDF title. Default font is Helvetica.""",
                },
                "content": {
                    "type": "string",
                    # 正文内容（纯文本，自动换行）。
                    "description": """Body content as plain text, wrapped automatically.""",
                },
                "sections": {
                    "type": "array",
                    # 章节列表。每项包含 heading 和 body。
                    "description": """List of sections. Each contains heading and body.""",
                    "items": {
                        "type": "object",
                        "properties": {
                            "heading": {"type": "string"},
                            "body": {"type": "string"},
                        },
                    },
                },
            },
        },
    },
    handler=_handle_write_pdf,
    emoji="📄",
    danger_level=ToolDangerLevel.write,
)